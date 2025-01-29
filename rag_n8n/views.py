from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

import io
import os

import pickle
from PyPDF2 import PdfReader
from pathlib import Path
from docx import Document
import pipmaster as pm

import logging

from .lightrag.lightrag import LightRAG, QueryParam
from .lightrag.llm import ollama_model_complete, ollama_embedding
from .lightrag.utils import EmbeddingFunc

from django.http import HttpResponse, StreamingHttpResponse

from rest_framework.decorators import authentication_classes, permission_classes
from adrf.decorators import api_view
from rest_framework.authentication import TokenAuthentication
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response

import aiofiles
from asgiref.sync import sync_to_async
import asyncio

from wakeonlan import send_magic_packet
import socket
import time


# Ollama server configuration
OLLAMA_HOST = "192.168.10.10" # Replace with your local IP of the computer serving Ollama
OLLAMA_PORT = 11434 # Port of Ollama
OLLAMA_MAC = "XX:XX:XX:XX:XX:XX"  # Replace with actual MAC address
WAKE_TIMEOUT = 60  # Seconds to wait for server to wake up


# Configure logging
logging.basicConfig(format="%(levelname)s:%(message)s", level=logging.INFO)


# Computer with Ollama serving wake up functions

async def wait_for_ollama_server():
    """Wait for Ollama server to become available"""
    start_time = time.time()
    while time.time() - start_time < WAKE_TIMEOUT:
        try:
            sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            sock.settimeout(2)
            result = sock.connect_ex((OLLAMA_HOST, OLLAMA_PORT))
            sock.close()
            if result == 0:
                logging.info("Ollama server is available")
                return True
        except Exception:
            pass
        await asyncio.sleep(2)
    return False

async def ensure_ollama_server():
    """Ensure Ollama server is running by sending WoL packet and waiting"""
    logging.info("Checking Ollama server availability...")
    
    # Try to connect first
    sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
    sock.settimeout(2)
    try:
        result = sock.connect_ex((OLLAMA_HOST, OLLAMA_PORT))
        if result == 0:
            logging.info("Ollama server is already running")
            sock.close()
            return True
    except Exception:
        pass
    finally:
        sock.close()

    # If server isn't responding, try to wake it
    logging.info("Sending Wake-on-LAN packet...")
    send_magic_packet(OLLAMA_MAC)
    
    # Wait for server to wake up
    if await wait_for_ollama_server():
        return True
    
    logging.error("Failed to wake Ollama server")
    return False


def start_LightRAG():
    logging.info("Initialize LightRAG...")
    WORKING_DIR = "./ollama"
    if not os.path.exists(WORKING_DIR):
        os.mkdir(WORKING_DIR)

    rag = LightRAG(
        working_dir=WORKING_DIR,
        llm_model_func=ollama_model_complete,
        llm_model_name="qwen2m",
        llm_model_max_async=4,
        llm_model_max_token_size=32768,
        llm_model_kwargs={"host": f"http://{OLLAMA_HOST}:{OLLAMA_PORT}", "options": {"num_ctx": 32768}},
        embedding_func=EmbeddingFunc(
            embedding_dim=768,
            max_token_size=8192,
            func=lambda texts: ollama_embedding(
                texts, embed_model="nomic-embed-text", host=f"http://{OLLAMA_HOST}:{OLLAMA_PORT}"
            ),
        ),
    )
    return rag

# Convert synchronous functions to async
start_LightRAG_async = sync_to_async(start_LightRAG)

async def get_drive_service_async():
    logging.info("Initializing Google Drive service...")
    SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
    creds = None
    
    if os.path.exists('token.pickle'):
        with open('token.pickle', 'rb') as token:
            creds = pickle.load(token)
    
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            await sync_to_async(creds.refresh)(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(
                'credentials.json', SCOPES)
            creds = await sync_to_async(flow.run_local_server)(port=0)
        with open('token.pickle', 'wb') as token:
            pickle.dump(creds, token)
    
    return await sync_to_async(build)('drive', 'v3', credentials=creds)

async def get_latest_document_async():
    logging.info("Fetching latest document...")
    service = await get_drive_service_async()
    query = "trashed=false"
    try:
        results = await sync_to_async(service.files().list)(
            q=query,
            orderBy="createdTime desc",
            pageSize=1,
            fields="files(id, name, createdTime, mimeType)"
        )
        results = await sync_to_async(results.execute)()
        
        files = results.get('files', [])
        return files[0] if files else None
    except Exception as e:
        logging.error(f"Error fetching document: {str(e)}")
        return None


async def download_and_process_file_async():
    logging.info("Processing document...")
    service = await get_drive_service_async()
    latest_file = await get_latest_document_async()
    
    if not latest_file:
        return "No files found in Google Drive."

    file_id = latest_file['id']
    file_name = latest_file['name']

    try:
        request = await sync_to_async(service.files().get_media)(fileId=file_id)
        file_content = io.BytesIO()
        downloader = MediaIoBaseDownload(file_content, request)
        
        done = False
        while not done:
            _, done = await sync_to_async(downloader.next_chunk)()
        
        file_content.seek(0)

        # File convertion to txt. In case handling more format is required try to use the textract python package. 
        match Path(file_name).suffix.lower():
            case ".pdf":
                pdf_reader = PdfReader(file_content)
                text = "\n".join(page.extract_text() for page in pdf_reader.pages)
            
            case ".txt" | ".md":
                text = file_content.getvalue().decode('utf-8')
                
            case ".docx":
                if not pm.is_installed("python-docx"):
                    pm.install("python-docx")
                doc = Document(file_content)
                text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
                
            case ".pptx":
                if not pm.is_installed("python-pptx"):
                    pm.install("python-pptx")
                from pptx import Presentation
                
                prs = Presentation(file_content)
                text = "\n".join(
                    shape.text 
                    for slide in prs.slides 
                    for shape in slide.shapes 
                    if hasattr(shape, "text")
                )
            
            case _:
                return "Error: Not compatible document format"
        
        output_filename = f"{os.path.splitext(file_name)[0]}.txt"
        async with aiofiles.open(output_filename, 'w', encoding='utf-8') as f:
            await f.write(text)
        return output_filename
            
    except Exception as e:
        return f"Error: {str(e)}"


@api_view(['POST'])
@authentication_classes([TokenAuthentication])
@permission_classes([IsAuthenticated])
async def process_data(request):
    logging.info('Processing RAG request...')
    data_mode = request.data.get('mode', 'mix')
    data_question = request.data.get('question')
    
    if not data_question:
        return Response({"error": "Invalid data"}, status=400)
    
    try:
        # Ensure Ollama server is running
        if not await ensure_ollama_server():
            return Response({
                "error": "Ollama server unavailable",
                "status": "error"
            }, status=503)
        
        output_text = await download_and_process_file_async()

        if "txt" not in output_text:
            return Response({"error": output_text}, status=400)
        
        # Initialize RAG
        rag = await start_LightRAG_async()
        
        # Read file content asynchronously
        async with aiofiles.open(f"./{output_text}", "r", encoding="utf-8") as f:
            content = await f.read()
            
            # Use event loop for insert operation
            loop = asyncio.get_event_loop()
            await loop.run_in_executor(None, lambda: rag.insert(content))
            
            # Use event loop for query operation
            response = await loop.run_in_executor(
                None,
                lambda: rag.query(
                    data_question,
                    param=QueryParam(
                        mode=data_mode,
                        only_need_context=False
                    )
                )
            )
            
            return Response({
                "status": "success",
                "response": response
            })

    except Exception as e:
        logging.error(f"Processing error: {str(e)}")
        return Response({
            "error": str(e),
            "status": "error"
        }, status=500)


def home(request):
    return HttpResponse("Welcome to the RAG API service")
